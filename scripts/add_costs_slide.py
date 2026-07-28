"""Met à jour la slide 8 (coûts d'exploitation) de Logiciel-douane_vF.pptx."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

PPTX_PATH = Path(r"c:\MAMP\htdocs\Logiciel-douane_vF.pptx")
PPTX_FALLBACK = Path(r"c:\MAMP\htdocs\Logiciel-douane_vF_jury.pptx")
SLIDE_INDEX = 10  # slide 11 (coûts d'exploitation dans le deck présenté)
ROI_SLIDE_INDEX = 8  # slide 9 — modèle visuel ROI
RESTORE_SLIDE8_EMPTY = False  # laisser la slide 8 inchangée si elle existe déjà

DARK = RGBColor(0x06, 0x2B, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x44, 0x37, 0x28)

# Taux indicatif pour les conversions affichées (1 USD ≈ 620 FCFA).
FCFA_PER_USD = 620

# Mesures terrain (juil. 2026, GPT-5, fiche formulaire complète).
BENCH_S500_USD = 0.07
BENCH_S500_SEC = 47
BENCH_RAMBO_RETRY_USD = 0.09
BENCH_RAMBO_SEC = 56
BENCH_RAMBO_FIRST_USD = 0.69  # pic observé au 1er essai (hors cache)
BENCH_DETAILED_AVG_USD = (BENCH_S500_USD + BENCH_RAMBO_RETRY_USD) / 2
MONTHLY_CLASSIFICATIONS = 100

PUNCHLINE = (
    "Une erreur de classification peut coûter plusieurs centaines de milliers de FCFA, "
    f"alors qu'une analyse Mosam coûte environ {int(BENCH_DETAILED_AVG_USD * FCFA_PER_USD)} FCFA "
    "(fiche détaillée mesurée)."
)

CALLOUT_BODY = (
    f"Mesures réelles : S500 neuf 2025 = {BENCH_S500_USD:.2f} $ ({BENCH_S500_SEC} s) ; "
    f"Rambo Magic = {BENCH_RAMBO_RETRY_USD:.2f} $ au 2e essai ({BENCH_RAMBO_SEC} s).\n"
    f"1er essai sans cache peut monter à ~{BENCH_RAMBO_FIRST_USD:.2f} $ "
    "(tokens de raisonnement GPT-5).\n"
    "Le cache après validation : requête identique à 0 $.\n"
    "Infrastructure cloud pilote : 0 à 30 $/mois. Moteur TEC / RGI : local, sans surcoût IA."
)

PIPELINE_BODY = (
    "Structure Pipeline\n\n"
    "1. Utilisateur (saisie)\n"
    "2. Identification GPT-5 (si fiche incomplète)\n"
    "3. Recherche web (si nécessaire)\n"
    "4. Embeddings OpenAI\n"
    "5. Index FAISS (TEC local)\n"
    "6. GPT-5 — positions candidates\n"
    "7. Moteur juridique Mosam\n"
    "   (Notes, RGI, Sous-positions, Droits)\n"
    "8. Résultat"
)


def _usd_to_fcfa(usd: float) -> int:
    return int(round(usd * FCFA_PER_USD))


def _fmt_usd_fr(usd: float, *, decimals: int = 2) -> str:
    return f"{usd:.{decimals}f}".replace(".", ",")


def _set_text(shape, text: str, *, font_name: str, size_pt: float, bold: bool = False, color=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _clone_text_style(source_shape, target_shape, text: str):
    src_p = source_shape.text_frame.paragraphs[0]
    src_run = src_p.runs[0] if src_p.runs else None
    tf = target_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if src_run is not None:
        run.font.name = src_run.font.name
        run.font.size = src_run.font.size
        run.font.bold = src_run.font.bold
        if src_run.font.color and src_run.font.color.type:
            run.font.color.rgb = src_run.font.color.rgb


def _add_punchline(cost_slide, prs: Presentation) -> None:
    ref_slide = prs.slides[4]
    ref_shape = next(
        s for s in ref_slide.shapes if hasattr(s, "text") and "Coût potentiel total" in s.text
    )
    box = cost_slide.shapes.add_textbox(
        ref_shape.left,
        4680000,
        ref_shape.width + 800000,
        ref_shape.height + 80000,
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = PUNCHLINE
    ref_run = ref_shape.text_frame.paragraphs[0].runs[0]
    run.font.name = ref_run.font.name
    run.font.size = ref_run.font.size
    run.font.bold = ref_run.font.bold
    run.font.color.rgb = ref_run.font.color.rgb


def _clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def main() -> None:
    target = PPTX_PATH if PPTX_PATH.exists() else PPTX_FALLBACK
    prs = Presentation(str(target))
    roi_slide = prs.slides[ROI_SLIDE_INDEX]
    cost_slide = prs.slides[SLIDE_INDEX]

    _clear_slide(cost_slide)

    if RESTORE_SLIDE8_EMPTY and len(prs.slides) > 7:
        _clear_slide(prs.slides[7])

    label_src = roi_slide.shapes[1]
    title_src = roi_slide.shapes[2]
    subtitle_src = roi_slide.shapes[3]
    num_src = roi_slide.shapes[5]
    item_title_src = roi_slide.shapes[6]
    item_body_src = roi_slide.shapes[7]
    callout_title_src = roi_slide.shapes[17]
    callout_body_src = roi_slide.shapes[18]
    callout_footer_src = roi_slide.shapes[20] if len(roi_slide.shapes) > 20 else roi_slide.shapes[19]

    monthly_usd = BENCH_DETAILED_AVG_USD * MONTHLY_CLASSIFICATIONS
    monthly_fcfa = _usd_to_fcfa(monthly_usd)

    def add_box(left, top, width, height):
        return cost_slide.shapes.add_textbox(left, top, width, height)

    label = add_box(label_src.left, label_src.top, label_src.width, label_src.height)
    _clone_text_style(label_src, label, "COÛTS D'EXPLOITATION")

    title = add_box(title_src.left, title_src.top, title_src.width, title_src.height)
    _clone_text_style(title_src, title, "Un modèle économique maîtrisé")

    subtitle = add_box(subtitle_src.left, subtitle_src.top, subtitle_src.width + 400000, subtitle_src.height)
    _clone_text_style(
        subtitle_src,
        subtitle,
        "Mesures réelles GPT-5 (juil. 2026) — fiche détaillée : ~0,07 à 0,10 $ / produit (~45–60 s)",
    )

    items = [
        (
            f"~{_fmt_usd_fr(BENCH_S500_USD)} $",
            "Fiche détaillée — Mercedes Classe S500",
            (
                f"Formulaire complet (essence, cylindrée, neuf…) : "
                f"{_fmt_usd_fr(BENCH_S500_USD)} $, {BENCH_S500_SEC} s "
                f"(~{_usd_to_fcfa(BENCH_S500_USD)} FCFA / produit)."
            ),
        ),
        (
            f"~{_fmt_usd_fr(BENCH_RAMBO_RETRY_USD)} $",
            "Fiche détaillée — Rambo Magic (2e essai)",
            (
                f"Même type de saisie, 2e classification : "
                f"{_fmt_usd_fr(BENCH_RAMBO_RETRY_USD)} $, {BENCH_RAMBO_SEC} s "
                f"(~{_usd_to_fcfa(BENCH_RAMBO_RETRY_USD)} FCFA). "
                f"1er essai observé : ~{_fmt_usd_fr(BENCH_RAMBO_FIRST_USD)} $."
            ),
        ),
        (
            f"~{_fmt_usd_fr(monthly_usd, decimals=0)} $/mois",
            f"Budget IA — {MONTHLY_CLASSIFICATIONS} classifications / mois",
            (
                f"Ordre de grandeur pilote (fiches détaillées, ~{_fmt_usd_fr(BENCH_DETAILED_AVG_USD)} $ en moyenne) : "
                f"~{monthly_fcfa} FCFA / mois IA "
                f"(~{_fmt_usd_fr(monthly_usd, decimals=0)} $/mois, hors cache)."
            ),
        ),
    ]

    num_positions = [roi_slide.shapes[5], roi_slide.shapes[9], roi_slide.shapes[13]]
    title_positions = [roi_slide.shapes[6], roi_slide.shapes[10], roi_slide.shapes[14]]
    body_positions = [roi_slide.shapes[7], roi_slide.shapes[11], roi_slide.shapes[15]]

    for idx, (num_text, item_title, item_body) in enumerate(items):
        nbox = add_box(
            num_positions[idx].left,
            num_positions[idx].top,
            num_positions[idx].width,
            num_positions[idx].height,
        )
        _clone_text_style(num_src, nbox, num_text)

        tbox = add_box(
            title_positions[idx].left,
            title_positions[idx].top,
            title_positions[idx].width,
            title_positions[idx].height,
        )
        _clone_text_style(item_title_src, tbox, item_title)

        bbox = add_box(
            body_positions[idx].left,
            body_positions[idx].top,
            body_positions[idx].width,
            body_positions[idx].height,
        )
        _clone_text_style(item_body_src, bbox, item_body)

    callout_bg = add_box(
        callout_title_src.left - 91440,
        callout_title_src.top - 91440,
        callout_title_src.width + 182880,
        callout_footer_src.top + callout_footer_src.height - callout_title_src.top + 182880,
    )
    callout_bg.fill.solid()
    callout_bg.fill.fore_color.rgb = DARK
    callout_bg.line.fill.background()

    ct = add_box(
        callout_title_src.left,
        callout_title_src.top,
        callout_title_src.width,
        callout_title_src.height,
    )
    _set_text(
        ct,
        "Coût marginal faible face au risque douanier",
        font_name="Crimson Pro Bold",
        size_pt=10,
        bold=True,
        color=WHITE,
    )

    cb = add_box(
        callout_body_src.left,
        callout_body_src.top,
        callout_body_src.width,
        callout_body_src.height + 280000,
    )
    _set_text(
        cb,
        CALLOUT_BODY,
        font_name="Open Sans",
        size_pt=8,
        color=WHITE,
    )

    pipeline = add_box(5160318, 1212800, 3600000, 3200000)
    _set_text(pipeline, PIPELINE_BODY, font_name="Open Sans", size_pt=9, color=MUTED)

    _add_punchline(cost_slide, prs)

    try:
        prs.save(str(target))
        print(f"Slide {SLIDE_INDEX + 1} mise à jour : {target}")
        if RESTORE_SLIDE8_EMPTY:
            print("Slide 8 vidée (coûts uniquement sur slide 11).")
    except PermissionError:
        prs.save(str(PPTX_FALLBACK))
        print(f"Fichier verrouillé — copie enregistrée : {PPTX_FALLBACK}")


if __name__ == "__main__":
    main()
